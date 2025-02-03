import fitz  # PyMuPDF
import io
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import traceback

def extract_images_and_text(pdf_file, pptx_file):
    try:
        print("Starting the conversion process...")  # Debugging start message
        
        # Open the PDF with PyMuPDF
        doc = fitz.open(pdf_file)
        prs = Presentation()  # Create a PowerPoint presentation object

        # Determine the directory to save temporary images (same directory as PDF or PowerPoint file)
        output_dir = os.path.dirname(pptx_file)  # Use the PowerPoint file's directory
        
        # Ensure directory exists
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Iterate over each page in the PDF
        for page_num in range(len(doc)):
            print(f"Processing page {page_num + 1}...")  # Debugging page number
            
            page = doc.load_page(page_num)
            
            # Create a new slide for each page of the PDF
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

            # Extract text from the page
            text = page.get_text("text")
            
            # Add text box to the slide, positioned at top-left of the slide
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
            text_frame = textbox.text_frame
            text_frame.text = text  # Set the extracted text

            # Extract images on the page
            img_list = page.get_images(full=True)
            for img_index, img in enumerate(img_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                # Convert image bytes to Image object (PIL)
                image = Image.open(io.BytesIO(image_bytes))
                
                # Save the image temporarily to insert into PowerPoint
                image_path = os.path.join(output_dir, f"temp_image_{page_num}_{img_index}.png")
                image.save(image_path)

                # Insert the image into the slide, positioned at the top-left corner
                slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(2.0))  # Adjust size as needed
                
                # Optionally remove the temporary image file
                os.remove(image_path)

        # Save the PowerPoint presentation
        prs.save(pptx_file)
        print(f"Conversion completed. The PowerPoint file is saved as {pptx_file}")
    
    except Exception as e:
        # If an error occurs, print the error message and log it
        print(f"ERROR: An error occurred during conversion: {e}")  # Print error to console
        
        # Ensure output_dir exists and is valid
        if not os.path.exists(output_dir):
            print("Error log directory does not exist. Creating directory.")
            os.makedirs(output_dir)
        
        # Log the error details to a file
        error_log_path = os.path.join(output_dir, "error_log.txt")
        try:
            with open(error_log_path, "w") as log_file:
                log_file.write(f"Error occurred: {str(e)}\n")
                log_file.write("Traceback:\n")
                traceback.print_exc(file=log_file)  # Log full traceback to error_log.txt
        except Exception as log_error:
            print(f"Error while writing to error log: {log_error}")
        
        # Print the path where the error log is saved
        print(f"An error occurred. Check the error log in: {error_log_path}")
        
        # Pauser to allow you to check the error log
        input("Press Enter to exit and check the error log...")

# Pauser to prevent automatic execution when you run the script
input("Press Enter to start the conversion process...")

# Example usage
pdf_file = r"C:\Your\File\Location\cv.pdf"  # Path to your PDF
pptx_file = r"C:\Your\File\Destination\output_with_images.pptx"  # Output PowerPoint file path

# Call the function to perform conversion
extract_images_and_text(pdf_file, pptx_file)

# Pauser to keep the program open and allow you to read any output or error messages
print("Conversion process completed or encountered an error.")
input("Press Enter to close the program after checking results...")
